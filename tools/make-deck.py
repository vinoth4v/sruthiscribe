#!/usr/bin/env python3
"""Build the SruthiScribe pitch deck.

    /tmp/pptxenv/bin/python tools/make-deck.py SruthiScribe.pptx

Every number in the deck comes from something that can be re-run -- the
accuracy figures from tools/eval.js, the library counts from the database, the
SSP yield from build-ssp.py --audit -- so the deck can be rebuilt rather than
edited when they change. Kept in the repo for that reason.
"""

import sys
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# The app's own palette, so the deck looks like the thing it describes.
INK      = RGBColor(0x1C, 0x19, 0x15)
INK_2    = RGBColor(0x45, 0x3F, 0x36)
MUTED    = RGBColor(0x75, 0x6D, 0x5E)
PAPER    = RGBColor(0xF6, 0xF3, 0xEC)
SURFACE  = RGBColor(0xFF, 0xFF, 0xFF)
ACCENT   = RGBColor(0xB3, 0x40, 0x1F)   # kumkum
TEAL     = RGBColor(0x0E, 0x6B, 0x5C)   # peacock
GOLD     = RGBColor(0x8A, 0x6D, 0x2F)
LINE     = RGBColor(0xE4, 0xDC, 0xC9)

W, H = Inches(13.333), Inches(7.5)
DISPLAY = "Georgia"
SANS = "Helvetica Neue"
MONO = "Menlo"


def deck():
    p = Presentation()
    p.slide_width, p.slide_height = W, H
    return p


def blank(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg = s.background.fill
    bg.solid()
    bg.fore_color.rgb = PAPER
    return s


def box(slide, x, y, w, h, fill=None, line=None):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    sh.shadow.inherit = False
    return sh


def text(slide, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         spacing=1.0):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    first = True
    for item in runs:
        s, size, colour, bold, font = (list(item) + [None] * 5)[:5]
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        para.line_spacing = spacing
        r = para.add_run()
        r.text = s
        r.font.size = Pt(size)
        r.font.color.rgb = colour or INK
        r.font.bold = bool(bold)
        r.font.name = font or SANS
    return tb


def title_slide(prs):
    s = blank(prs)
    slab = s.shapes.add_shape(1, 0, 0, Inches(0.16), H)   # rectangle
    slab.fill.solid(); slab.fill.fore_color.rgb = ACCENT
    slab.line.fill.background(); slab.shadow.inherit = False
    text(s, Inches(1.1), Inches(2.3), Inches(11), Inches(1.1),
         [("SruthiScribe", 60, INK, True, DISPLAY)])
    text(s, Inches(1.1), Inches(3.35), Inches(10.5), Inches(0.9),
         [("Sing it. Read it back as svaras.", 26, ACCENT, False, DISPLAY)])
    text(s, Inches(1.1), Inches(4.25), Inches(10.2), Inches(1.4),
         [("A browser-based Carnatic transcription engine, an open machine-readable "
           "library of the repertoire, and guided practice against it.",
           16, INK_2, False, SANS)])
    text(s, Inches(1.1), Inches(6.02), Inches(11), Inches(0.95), [
        ("Vinoth Kannan  ·  vinoth.kannan.eu@gmail.com", 12.5, GOLD, False, MONO),
        ("Beta  ·  measured accuracy, published limits  ·  August 2026",
         12, MUTED, False, MONO),
    ], spacing=1.5)
    return s


def section(prs, kicker, heading):
    s = blank(prs)
    text(s, Inches(1.1), Inches(3.0), Inches(11), Inches(0.4),
         [(kicker.upper(), 13, ACCENT, True, MONO)])
    text(s, Inches(1.1), Inches(3.5), Inches(11), Inches(1.0),
         [(heading, 40, INK, True, DISPLAY)])
    return s


def content(prs, kicker, heading):
    s = blank(prs)
    text(s, Inches(0.9), Inches(0.55), Inches(11.5), Inches(0.3),
         [(kicker.upper(), 11.5, ACCENT, True, MONO)])
    text(s, Inches(0.9), Inches(0.95), Inches(11.5), Inches(0.7),
         [(heading, 30, INK, True, DISPLAY)])
    return s


def bullets(slide, x, y, w, items, size=15.5, gap=0.52):
    for i, (head, body) in enumerate(items):
        yy = y + Inches(i * gap)
        text(slide, x, yy, w, Inches(0.4),
             [(head, size, INK, True, SANS)])
        if body:
            text(slide, x, yy + Inches(0.3), w, Inches(0.8),
                 [(body, size - 2.5, MUTED, False, SANS)], spacing=1.25)


def stat(slide, x, y, w, big, label, colour=ACCENT):
    box(slide, x, y, w, Inches(1.5), SURFACE, LINE)
    text(slide, x, y + Inches(0.24), w, Inches(0.62),
         [(big, 34, colour, True, DISPLAY)], align=PP_ALIGN.CENTER)
    text(slide, x, y + Inches(0.92), w, Inches(0.5),
         [(label, 11, MUTED, False, SANS)], align=PP_ALIGN.CENTER)


def build(path):
    prs = deck()
    title_slide(prs)

    # ---------------------------------------------------------------- problem
    s = content(prs, "The problem", "Carnatic music is taught by ear, and stays there")
    bullets(s, Inches(0.9), Inches(1.95), Inches(5.6), [
        ("A student cannot see what they just sang",
         "Correction happens in the lesson, from the teacher's ear. Between lessons "
         "there is no feedback at all."),
        ("Notation exists, but not where you need it",
         "Printed books, PDFs and scanned archives — none of it queryable, none of it "
         "aligned to what you are singing."),
        ("The reference corpus is scattered and closed",
         "Ragam definitions on one site, lyrics on another, notation in a 1904 book. "
         "Little of it is machine-readable."),
    ], gap=1.06)
    box(s, Inches(7.0), Inches(1.95), Inches(5.4), Inches(3.2), SURFACE, LINE)
    text(s, Inches(7.35), Inches(2.25), Inches(4.7), Inches(2.6), [
        ("What a learner actually has today", 15, INK, True, SANS),
        ("", 8, MUTED, False, SANS),
        ("A tanpura app for the drone.", 13, INK_2, False, SANS),
        ("A tuner that says whether one note is sharp.", 13, INK_2, False, SANS),
        ("A PDF of notation they cannot search.", 13, INK_2, False, SANS),
        ("A recording of their own practice they will never listen back to.",
         13, INK_2, False, SANS),
        ("", 8, MUTED, False, SANS),
        ("Nothing that turns the singing into notation.", 13, ACCENT, True, SANS),
    ], spacing=1.35)

    # ------------------------------------------------------------------- need
    s = content(prs, "The need", "A loop, not a tool")
    panels = [
        (ACCENT, "1 \u00b7 Hear yourself",
         "Sing a phrase; get back svaras with sthayi, gamaka and avartana bars, "
         "laid out the way the tradition prints it."),
        (TEAL, "2 \u00b7 Find it written",
         "Compositions, ragams, talams and notation in one place, searchable, with "
         "the source and licence of every row recorded."),
        (GOLD, "3 \u00b7 Practise against it",
         "Sing a written composition and be told, svara by svara, which you held "
         "and which you missed \u2014 in cents, not opinion."),
    ]
    for i, (colour, head, body) in enumerate(panels):
        x = Inches(0.9 + i * 3.94)
        box(s, x, Inches(1.9), Inches(3.64), Inches(2.7), SURFACE, LINE)
        text(s, Inches(0.9 + i * 3.94 + 0.3), Inches(2.2), Inches(3.04), Inches(2.3), [
            (head, 15.5, colour, True, DISPLAY),
            ("", 9, MUTED, False, SANS),
            (body, 12.5, INK_2, False, SANS),
        ], spacing=1.3)
    text(s, Inches(0.9), Inches(4.95), Inches(11.5), Inches(0.8),
         [("Each stage feeds the next. The library gives the decoder its ragam and tala; "
           "the decoder gives the library new readings; the notation gives practice "
           "something exact to score against. A student can sing a kriti they cannot yet "
           "write, contribute it, and then practise the thing they just wrote.",
           13, MUTED, False, SANS)], spacing=1.25)

    # ---------------------------------------------------------- functionality
    s = content(prs, "Functionality", "What it does today")
    bullets(s, Inches(0.9), Inches(2.05), Inches(5.5), [
        ("Record or upload, and read it back",
         "Microphone or file, in the browser. Svaras with octave, duration, gamaka "
         "class and cents deviation."),
        ("Tala-aware layout",
         "All seven suladi sapta talas × five jatis, nadai 1x–8x, avartana and anga "
         "bar lines, sahitya under the svaras."),
        ("Fix what it got wrong",
         "Edit any svara, mark the eduppu, regroup speeds — the decode is a starting "
         "point, not a verdict."),
        ("Export",
         "Print-quality PDF laid out as the tradition prints it, JSON, or contribute "
         "the reading to the library."),
    ], gap=1.06)
    bullets(s, Inches(6.9), Inches(2.05), Inches(5.5), [
        ("A library of 598 compositions",
         "42 with complete notation, searchable by ragam, tala, composer, deity and "
         "how complete the notation is."),
        ("939 ragams",
         "All 72 melakartas and their janyas, with arohana/avarohana; 120 wired into "
         "the decoder as constraints."),
        ("Provenance on every row",
         "Source, licence and page, so a notation can be checked against the book it "
         "came from."),
        ("Sadhana \u2014 guided practice",
         "Sing any notated composition against a scrolling stage with a drone and a "
         "tala click; scored per svara in cents."),
    ], gap=1.06)

    # ------------------------------------------------------------ how it works
    s = content(prs, "Technical architecture", "A signal-processing pipeline, entirely in the browser")
    stages = [
        ("Audio", "mic or file\n16 kHz mono"),
        ("YIN", "pitch track\n50 ms window"),
        ("Voicing", "energy + confidence\n90 ms min. note"),
        ("Tonic", "pitch histogram\nauto-sruthi"),
        ("Viterbi", "ragam-constrained\nsvara decode"),
        ("Notation", "tala grid\nPDF / JSON"),
    ]
    # Derive the box width from the space available rather than a guessed
    # constant: six at 1.83 with 0.28 gaps ended 0.05in from the slide edge and
    # the last one was clipped.
    gap = Inches(0.24)
    wbox = int((W - Inches(1.8) - gap * (len(stages) - 1)) / len(stages))
    x = Inches(0.9)
    for i, (head, body) in enumerate(stages):
        colour = TEAL if i < 3 else ACCENT
        box(s, x, Inches(2.0), wbox, Inches(1.55), SURFACE, LINE)
        text(s, x, Inches(2.22), wbox, Inches(0.4),
             [(head, 15, colour, True, SANS)], align=PP_ALIGN.CENTER)
        text(s, x, Inches(2.66), wbox, Inches(0.8),
             [(body, 10.5, MUTED, False, MONO)], align=PP_ALIGN.CENTER, spacing=1.25)
        if i < len(stages) - 1:
            text(s, x + wbox, Inches(2.55), Inches(0.28), Inches(0.4),
                 [("›", 20, LINE, True, SANS)], align=PP_ALIGN.CENTER)
        x += wbox + Inches(0.28)

    bullets(s, Inches(0.9), Inches(4.1), Inches(5.5), [
        ("No server, no model, no key",
         "The engine is 46 kB of dependency-free JavaScript. Transcription works "
         "offline; nothing is uploaded."),
        ("The ragam is the decoder's state space",
         "Viterbi over that ragam's svarasthanas, penalising switches — so the decode "
         "cannot invent a note outside the ragam."),
    ], gap=0.9)
    bullets(s, Inches(6.9), Inches(4.1), Inches(5.5), [
        ("Postgres + REST for the library",
         "Supabase; the page talks to it directly. Notation stored as JSON with the "
         "tala and akshara grid."),
        ("The harness scores the shipping code",
         "tools/eval.js extracts the engine from the page itself and scores it against "
         "public datasets — the two cannot drift."),
    ], gap=0.9)

    # -------------------------------------------------------------- measured
    s = content(prs, "Measured, not claimed", "Accuracy on studio recordings")
    stat(s, Inches(0.9), Inches(1.95), Inches(2.7), "91.7%", "svara accuracy\n(label + octave)")
    stat(s, Inches(3.85), Inches(1.95), Inches(2.7), "82.5%", "raw pitch accuracy", TEAL)
    stat(s, Inches(6.8), Inches(1.95), Inches(2.7), "2.5%", "octave error rate", TEAL)
    stat(s, Inches(9.75), Inches(1.95), Inches(2.65), "3.2¢", "tonic (sruthi) error", TEAL)
    text(s, Inches(0.9), Inches(3.75), Inches(11.5), Inches(0.5),
         [("Scored on 33 items of Sanidha (Georgia Tech, CC BY 4.0), studio multitrack, "
           "against the dataset's own pitch annotations. It declines 27% of sung frames "
           "rather than guess them.", 13, MUTED, False, SANS)])
    box(s, Inches(0.9), Inches(4.35), Inches(11.5), Inches(1.8), SURFACE, LINE)
    text(s, Inches(1.25), Inches(4.6), Inches(10.8), Inches(1.7), [
        ("How the measurement is kept honest", 15, INK, True, SANS),
        ("The harness runs the engine that ships, not a copy. Every change must improve "
         "two datasets (Sanidha and Saraga), both halves of a holdout split, and survive "
         "a deliberately mis-set tonic before it lands — what helps one corpus and "
         "hurts the other is rejected. Annotation faults are handled in the open: three "
         "records whose contour tracks another instrument are excluded, and frames the "
         "dataset itself doubled an octave are repaired by rules that read only the "
         "reference, never the engine's output. The web page publishes these same "
         "figures, and a test fails if they drift from the measurement.",
         12.5, INK_2, False, SANS),
    ], spacing=1.3)

    # ------------------------------------------------------------------- SSP
    s = content(prs, "The library", "Rescuing a 1904 book by machine")
    text(s, Inches(0.9), Inches(1.9), Inches(11.4), Inches(1.0), [
        ("Subbarama Dikshitar's Sangita Sampradaya Pradarsini (1904) is the only "
         "public-domain source of complete Carnatic notation — and it has never been "
         "machine-readable.", 15, INK_2, False, SANS)], spacing=1.3)
    stat(s, Inches(0.9), Inches(3.0), Inches(2.7), "598", "compositions", TEAL)
    stat(s, Inches(3.85), Inches(3.0), Inches(2.7), "42", "with complete notation")
    stat(s, Inches(6.8), Inches(3.0), Inches(2.7), "939", "ragams", TEAL)
    stat(s, Inches(9.75), Inches(3.0), Inches(2.65), "4,695", "svaras extracted from SSP")
    text(s, Inches(0.9), Inches(4.85), Inches(11.4), Inches(1.6), [
        ("The extractor reads the book's own notation scheme from its page iii — "
         "capitals are two aksharas, underlines halve, dots mark sthayi — and verifies "
         "every section against the tala the book prints for it. 26 kirtanas parse "
         "clean end to end; 12 of them were confirmed independently against a separate "
         "catalogue, agreeing on both title and ragam. Nothing that fails verification "
         "is loaded.", 13.5, INK_2, False, SANS)], spacing=1.35)

    # ----------------------------------------------------------- limitations
    s = content(prs, "Limitations", "What it cannot do yet")
    bullets(s, Inches(0.9), Inches(2.05), Inches(5.5), [
        ("Ragam identification is weak — 15% top-1",
         "A pitch histogram cannot separate ragams that share a svara set and differ "
         "by phrase. The user still has to name the ragam."),
        ("Gamaka classification is coarse",
         "Five machine labels, not the dasavidha taxonomy. It says 'kampita', not which "
         "kampita."),
        ("It abstains, and it false-alarms",
         "27% of sung frames get no svara — declined, not guessed — and 32% of "
         "silences still get called a note."),
        ("One voice, no ensemble",
         "Trained on and scored against solo vocal. Accompaniment in the room degrades it."),
    ], gap=1.06)
    bullets(s, Inches(6.9), Inches(2.05), Inches(5.5), [
        ("The library is mostly metadata",
         "433 of 598 rows have no notation. Complete notation is not openly licensed at "
         "scale; that is a licensing wall, not an engineering one."),
        ("The SSP extractor verifies 44% of sections",
         "Errors are per-glyph now — underline geometry, tight spacing. Each further "
         "point is slow work."),
        ("No Tyagaraja notation",
         "SSP covers Dikshitar. The Walajapet manuscripts are scans with no text layer."),
        ("Not a substitute for a teacher",
         "It measures pitch. It does not hear bhava, and it should not be sold as if "
         "it did."),
    ], gap=1.06)

    # ------------------------------------------------------------- landscape
    # Each entry gets a marker and a gap: run together as plain paragraphs they
    # read as one block of prose and the list stops being a list.
    def listed(items, colour=INK_2, bold=False):
        out = []
        for i, t in enumerate(items):
            out.append(("\u00b7  " + t, 12.5, colour, bold, SANS))
            if i < len(items) - 1:
                out.append(("", 6, MUTED, False, SANS))
        return out

    s = content(prs, "Landscape", "What is genuinely new here")
    box(s, Inches(0.9), Inches(1.85), Inches(5.6), Inches(3.6), SURFACE, LINE)
    text(s, Inches(1.25), Inches(2.15), Inches(4.9), Inches(3.2),
         [("Already exists", 16, MUTED, True, DISPLAY), ("", 9, MUTED, False, SANS)] +
         listed([
             "Tuners that name the svara you are holding \u2014 Carnatic Tuner, Shruti.",
             "Practice apps with pitch feedback and lesson tracks \u2014 Riyaz, "
             "Carnatic Singer. Sadhana is our answer to these, and they got there first.",
             "AI accompaniment and raga detection \u2014 NaadSadhana, an Apple Design "
             "Award winner with 410 ragas.",
             "Notation typing tools \u2014 Haathi Carnatic Editor.",
             "Research corpora \u2014 CompMusic, Saraga, Sanidha.",
         ]), spacing=1.25)

    box(s, Inches(6.9), Inches(1.85), Inches(5.5), Inches(3.6), SURFACE, LINE)
    text(s, Inches(7.25), Inches(2.15), Inches(4.8), Inches(3.2),
         [("Less crowded", 16, ACCENT, True, DISPLAY), ("", 9, MUTED, False, SANS)] +
         listed([
             "Continuous transcription of a phrase into tala-aligned notation, not "
             "just the current note.",
             "Runs entirely in the browser \u2014 no install, no account, works offline.",
             "Published accuracy against public datasets, with the harness in the repo.",
             "Practice scored against an open library rather than a proprietary "
             "lesson set \u2014 including a 1904 book nobody else can drill you on.",
         ]) +
         [("", 6, MUTED, False, SANS)] +
         listed(["A machine-readable SSP corpus with provenance \u2014 as far as we "
                 "know, the first."], ACCENT, True), spacing=1.25)

    text(s, Inches(0.9), Inches(5.68), Inches(11.5), Inches(0.8),
         [("Honest reading: transcription is differentiated, not unique, and practice "
           "apps had a decade\u2019s head start. What none of them close is the loop \u2014 "
           "sing something down, contribute it, then be drilled on it. The extracted "
           "SSP corpus is what makes that loop worth closing.",
           12.5, MUTED, False, SANS)], spacing=1.25)

    # ------------------------------------------------------------------ close
    s = section(prs, "Where it goes", "Next")
    text(s, Inches(1.1), Inches(4.55), Inches(11), Inches(2.0), [
        (t, 15, c, b, f) for (t, _sz, c, b, f) in listed([
            "Raise the SSP yield \u2014 every point converts directly into complete works.",
            "Cut the voicing false alarm; it is the weakest measured number.",
            "Phrase-level ragam identification, where the published work actually lives.",
            "Attach openly-licensed audio to notated works, so a student can hear and "
            "read together.",
            "Take Sadhana from pitch to timing \u2014 holding a svara for the right "
            "number of aksharas is half of singing in tala.",
        ])], spacing=1.25)
    text(s, Inches(1.1), Inches(6.98), Inches(11), Inches(0.42),
         [("Questions, corrections, missing compositions \u2014 "
           "Vinoth Kannan, vinoth.kannan.eu@gmail.com  \u00b7  \u00a9 2026, all rights "
           "reserved", 12.5, MUTED, False, MONO)])

    prs.save(path)
    return len(prs.slides.__iter__.__self__._sldIdLst)


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "SruthiScribe.pptx"
    n = build(out)
    print("wrote %s (%d slides)" % (out, n))
